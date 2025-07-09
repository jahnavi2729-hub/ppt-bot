import os, uuid, wikipedia
from flask import Flask, request, jsonify, render_template, send_from_directory
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
import requests

# Load .env variables
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

# Flask app setup
app = Flask(__name__, static_folder="static")
OUTPUT_DIR = os.path.join(app.static_folder, "ppt")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Wikipedia summary
def get_wikipedia_summary(topic: str) -> str:
    try:
        return wikipedia.summary(topic, sentences=5)
    except Exception as e:
        return f"(Wikipedia retrieval failed: {e})"

# Prompt builder
def build_prompt(topic: str, context: str) -> str:
    return (
        f"Using the context below, generate only bullet-point content to fill 4-5 PowerPoint slides on '{topic}'. "
        f"Each slide should have 2-3 bullet points. Do NOT include a title slide or thank-you slide. Just bullet points only.\n\n"
        f"Context:\n{context}\n\nSlides:"
    )

# OpenRouter API call
def call_openrouter_llm(prompt: str) -> str:
    try:
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "HTTP-Referer": "http://localhost:5000",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "mistralai/mistral-7b-instruct",
            "messages": [{"role": "user", "content": prompt}]
        }
        response = requests.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
        if response.status_code == 200:
            return_value = response.json()["choices"][0]["message"]["content"]
            print(return_value)
            return return_value

        return f"LLM Error: {response.status_code} - {response.text}"
    except Exception as e:
        return f"LLM Error: {e}"

# Group list helper
def chunk_list(lst, size):
    return [lst[i:i+size] for i in range(0, len(lst), size)]

# PPT generator
def generate_ppt(topic: str, slide_text: str) -> str:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    import re
    import uuid
    import os

    prs = Presentation()

    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = topic.title()
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # Clean bullet text
    lines = [line.strip() for line in slide_text.split('\n') if line.strip()]
    bullet_lines = []
    for line in lines:
        if line.lower().startswith("slide ") and ":" in line:
            continue
        line = re.sub(r"^(?:\d+\.\s*|[-•]\s*)", "-", line).strip()
        bullet_lines.append(line)

    # Group into 2–3 bullet points per slide
    grouped = [bullet_lines[i:i + 3] for i in range(0, len(bullet_lines), 3)][:5]

    # Content Slides with proper bullet formatting
    for group in grouped:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()  # Clear default paragraph

        for bullet in group:
            para = tf.add_paragraph()
            para.text = bullet
            para.level = 0
            para.bullet = True  # ✅ This enables bullet point
            para.font.size = Pt(24)
            para.font.name = 'Calibri'
            para.font.color.rgb = RGBColor(0, 0, 0)

    # Thank You Slide
    thank_slide = prs.slides.add_slide(prs.slide_layouts[6])
    thank_box = thank_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = thank_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 102, 0)
    p.alignment = PP_ALIGN.CENTER

    # Save PPT file
    fname = f"{topic.replace(' ', '_')[:20]}_{uuid.uuid4().hex[:5]}.pptx"
    path = os.path.join(OUTPUT_DIR, fname)
    prs.save(path)
    return fname
# Routes
@app.route("/")
def home():
    return render_template("chat.html")

@app.route("/ask", methods=["POST"])
def ask():
    topic = (request.json or {}).get("message", "").strip()
    if not topic:
        return jsonify(answer="❓ Please enter a topic.")

    wiki_text = get_wikipedia_summary(topic)
    prompt = build_prompt(topic, wiki_text)
    response = call_openrouter_llm(prompt)

    if response.startswith("LLM Error"):
        return jsonify(answer=response)

    ppt_file = generate_ppt(topic, response)
    return jsonify(answer="✅ PPT generated!", ppt=f"/static/ppt/{ppt_file}")

@app.route("/static/ppt/<path:filename>")
def download(filename):
    return send_from_directory(OUTPUT_DIR, filename, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True, port=5000)
